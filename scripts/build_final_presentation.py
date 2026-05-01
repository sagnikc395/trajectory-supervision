from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "final" / "final_presentation.pptx"
FIG_DIR = ROOT / "CS_590NN_690NN_Final_Project_Template" / "figures"

WIDE_W = Inches(13.333)
WIDE_H = Inches(7.5)

INK = RGBColor(28, 38, 52)
MUTED = RGBColor(89, 104, 128)
RULE = RGBColor(215, 220, 228)
PALE = RGBColor(246, 248, 251)
TEAL = RGBColor(18, 132, 121)
RED = RGBColor(181, 71, 69)
GOLD = RGBColor(174, 122, 37)
WHITE = RGBColor(255, 255, 255)

FONT = "Helvetica Neue"


def add_text(slide, text, x, y, w, h, size=20, bold=False, color=INK,
             align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.03)
    frame.margin_right = Inches(0.03)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.alignment = align
    p.space_after = Pt(3)
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_lines(slide, lines, x, y, w, h, size=18, color=INK, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    for idx, line in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = line
        p.font.name = FONT
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.space_after = Pt(gap)
    return box


def add_rule(slide, y):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.72), Inches(y), Inches(11.9), Inches(0.02)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RULE
    shape.line.fill.background()


def add_title(slide, title, subtitle=None):
    add_text(slide, title, 0.72, 0.38, 11.8, 0.43, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.74, 0.86, 11.4, 0.28, size=11, color=MUTED)
    add_rule(slide, 1.18)


def add_footer(slide, idx):
    add_text(slide, "CS 590NN Final Project", 0.72, 7.07, 3.0, 0.18, size=7, color=MUTED)
    add_text(slide, str(idx), 12.38, 7.07, 0.25, 0.18, size=7, color=MUTED, align=PP_ALIGN.RIGHT)


def add_pill(slide, label, x, y, w, h, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.8)
    add_text(slide, label, x + 0.12, y + 0.12, w - 0.24, h - 0.18, size=17, bold=True,
             color=color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def add_table(slide, rows, x, y, w, h, widths=None, font_size=12):
    table_shape = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h))
    table = table_shape.table
    if widths:
        for i, width in enumerate(widths):
            table.columns[i].width = Inches(width)
    for r, row in enumerate(rows):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.margin_left = Inches(0.06)
            cell.margin_right = Inches(0.06)
            cell.margin_top = Inches(0.03)
            cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            cell.fill.fore_color.rgb = INK if r == 0 else (PALE if r % 2 else WHITE)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT
            p.font.size = Pt(font_size if r else font_size - 1)
            p.font.bold = r == 0
            p.font.color.rgb = WHITE if r == 0 else INK
    return table_shape


def add_image_fit(slide, path, x, y, w, h):
    pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w))
    scale = min(Inches(w) / pic.width, Inches(h) / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = Inches(x) + int((Inches(w) - pic.width) / 2)
    pic.top = Inches(y) + int((Inches(h) - pic.height) / 2)
    return pic


def add_stat(slide, value, label, x, y, w, color):
    add_text(slide, value, x, y, w, 0.48, size=29, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x, y + 0.52, w, 0.24, size=9, color=MUTED, align=PP_ALIGN.CENTER)


def build():
    prs = Presentation()
    prs.slide_width = WIDE_W
    prs.slide_height = WIDE_H
    blank = prs.slide_layouts[6]

    # 1. Cover
    slide = prs.slides.add_slide(blank)
    add_text(slide, "Trajectory supervision for\ncontinual tool-use learning",
             0.78, 0.82, 7.2, 1.2, size=31, bold=True)
    add_text(slide, "Single-seed Llama 3.1 8B QLoRA pilot on API-Bank",
             0.82, 2.23, 7.2, 0.32, size=15, color=MUTED)
    add_text(slide, "+17.7 pts", 9.08, 1.02, 2.6, 0.55, size=31, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(slide, "final exact full-call accuracy\nB over A",
             9.04, 1.78, 2.7, 0.42, size=12, color=INK, align=PP_ALIGN.CENTER)
    add_rule(slide, 3.02)
    add_lines(slide, [
        "Question: does API trajectory context help continual tool-use fine-tuning?",
        "Result: B improves exact calls and API selection.",
        "Caveat: one seed, and B has 25.1% more training tokens."
    ], 0.88, 3.45, 10.7, 1.55, size=18, gap=10)
    add_text(slide, "Vishnu Vardhan Reddy B, Sagnik Chatterjee, Soumik Bhatta",
             0.82, 6.08, 9.0, 0.25, size=12)
    add_text(slide, "CS 590NN / 690NN Final Project", 0.82, 6.38, 4.0, 0.2, size=9, color=MUTED)
    add_footer(slide, 1)

    # 2. Motivation
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Motivation", "The useful signal is the observable tool-use trace.")
    add_text(slide, "Final answers hide the path", 0.95, 1.72, 5.4, 0.38, size=22, bold=True)
    add_lines(slide, [
        "Tool-use examples include API calls, arguments, and observations.",
        "Those traces are not private reasoning.",
        "They still give an action-observation history before the next API call."
    ], 0.98, 2.72, 5.5, 1.7, size=17, gap=10)
    add_text(slide, "Our comparison", 7.18, 1.72, 4.0, 0.34, size=22, bold=True, color=TEAL)
    add_lines(slide, [
        "A: stripped context",
        "B: trajectory context",
        "Same base model, seed, stream order, and scoring code."
    ], 7.22, 2.35, 4.7, 1.4, size=18, gap=12)
    add_text(slide, "The question is whether keeping the trace changes continual next-call learning.",
             1.0, 5.4, 11.1, 0.45, size=20, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # 3. Experiment design
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Experiment Design", "Train sequentially; evaluate every block after each stage.")
    x0 = 0.96
    for i, label in enumerate(["D1", "D2", "D3", "D4"]):
        x = x0 + i * 2.75
        add_pill(slide, label, x, 2.08, 1.35, 0.7, TEAL)
        add_text(slide, "train", x, 2.86, 1.35, 0.18, size=8, color=MUTED, align=PP_ALIGN.CENTER)
        if i < 3:
            add_text(slide, ">", x + 1.58, 2.25, 0.45, 0.24, size=19, bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, "After each training stage: score D1, D2, D3, and D4",
             1.15, 3.62, 10.7, 0.34, size=20, bold=True, align=PP_ALIGN.CENTER)
    add_rule(slide, 4.35)
    add_text(slide, "Condition A", 1.05, 5.02, 2.0, 0.3, size=20, bold=True, color=RED)
    add_text(slide, "Stripped-context next-API-call baseline.", 3.2, 5.06, 4.3, 0.24, size=15)
    add_text(slide, "Condition B", 1.05, 5.72, 2.0, 0.3, size=20, bold=True, color=TEAL)
    add_text(slide, "Trajectory-context next-API-call format.", 3.2, 5.76, 4.5, 0.24, size=15)
    add_footer(slide, 3)

    # 4. Setup
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Setup", "Same model and scoring pipeline; B sees more prompt tokens.")
    rows = [
        ["Component", "Value"],
        ["Model", "Llama 3.1 8B Instruct"],
        ["Fine-tuning", "QLoRA, rank 32, alpha 64, 4-bit"],
        ["Data stream", "API-Bank split into D1-D4"],
        ["Seed", "42"],
        ["Full eval examples", "126, 104, 103, 107"],
        ["Train tokens", "A 1.86M, B 2.32M (+25.1%)"],
    ]
    add_table(slide, rows, 1.0, 1.55, 7.15, 3.42, widths=[2.15, 5.0], font_size=11)
    add_text(slide, "Scoring", 9.02, 1.72, 2.3, 0.32, size=20, bold=True, color=TEAL)
    add_lines(slide, [
        "API-name accuracy",
        "Exact full-call accuracy",
        "Name + any parameter",
        "Malformed or no-call rate"
    ], 9.04, 2.28, 3.25, 1.55, size=15, gap=8)
    add_text(slide, "Token count is the main fairness caveat.", 1.05, 5.72, 10.7, 0.32,
             size=19, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # 5. Main result
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Main Result", "B has higher final exact full-call accuracy on every held-out block.")
    add_stat(slide, "39.2%", "A mean exact", 1.02, 1.42, 1.7, RED)
    add_stat(slide, "56.9%", "B mean exact", 3.0, 1.42, 1.7, TEAL)
    add_stat(slide, "+17.7", "point gap", 4.98, 1.42, 1.7, GOLD)
    chart_data = CategoryChartData()
    chart_data.categories = ["D1", "D2", "D3", "D4"]
    chart_data.add_series("Condition A", [35.7, 43.3, 32.0, 45.8])
    chart_data.add_series("Condition B", [57.9, 61.5, 44.7, 63.6])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0),
        Inches(2.52),
        Inches(10.9),
        Inches(3.55),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 80
    for series, color in zip(chart.series, [RED, TEAL]):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = color
        series.format.line.color.rgb = color
    add_text(slide, "Still not causal: B also trains on more tokens.", 7.72, 1.58, 4.3, 0.28,
             size=15, color=MUTED, align=PP_ALIGN.RIGHT)
    add_footer(slide, 5)

    # 6. Continual matrix
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Continual Evaluation", "We score all blocks after each training stage.")
    add_image_fit(slide, FIG_DIR / "full_eval_exact_full_heatmaps_seed42.png", 1.05, 1.55, 10.5, 4.7)
    add_text(slide, "B is higher on every final-stage block and most earlier train/eval pairs.",
             1.2, 6.28, 10.4, 0.3, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    # 7. Failure analysis
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Failure Analysis", "Trajectory context improves API choice but increases parse failures.")
    rows = [
        ["Category", "A", "B", "Read"],
        ["Exact full calls", "172", "251", "B higher"],
        ["Wrong API", "102", "12", "B lower"],
        ["Wrong / partial params", "121", "76", "B lower"],
        ["Malformed / no call", "45", "101", "B higher"],
    ]
    add_table(slide, rows, 0.9, 1.58, 11.35, 2.7, widths=[4.0, 1.1, 1.1, 5.15], font_size=12)
    add_text(slide, "Takeaway", 1.0, 5.08, 2.0, 0.3, size=21, bold=True, color=TEAL)
    add_text(slide, "B chooses the right API more often, but exact formatting is less stable.",
             3.02, 5.1, 8.8, 0.32, size=18)
    add_footer(slide, 7)

    # 8. Conclusion
    slide = prs.slides.add_slide(blank)
    add_title(slide, "Takeaways", "Current evidence is a single-seed pilot.")
    add_text(slide, "What we found", 0.95, 1.65, 3.4, 0.35, size=22, bold=True, color=TEAL)
    add_lines(slide, [
        "Trajectory context favored exact calls: 56.9% vs 39.2%.",
        "It strongly reduced wrong-API errors: 12 vs 102."
    ], 0.98, 2.28, 4.8, 1.1, size=16, gap=9)
    add_text(slide, "Limits", 0.95, 4.05, 2.2, 0.35, size=22, bold=True, color=RED)
    add_lines(slide, [
        "One seed.",
        "B has more tokens.",
        "Next-call benchmark, not full task success."
    ], 0.98, 4.66, 4.8, 1.25, size=16, gap=9)
    add_text(slide, "Next", 7.3, 1.65, 2.0, 0.35, size=22, bold=True, color=GOLD)
    add_lines(slide, [
        "Run multiple seeds.",
        "Train a token-matched B condition.",
        "Add semantic argument scoring."
    ], 7.32, 2.28, 4.6, 1.3, size=17, gap=10)
    add_text(slide, "Bottom line: trajectory context helps tool identity, but exact call formatting still needs work.",
             1.0, 6.35, 11.4, 0.35, size=18, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 8)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
